import os
import platform
import re
import time
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from .layout import add_images_to_slide, get_resource_path
from .image_formats import is_image_file, register_heic_support

register_heic_support()
DATETIME_FMT = "%Y-%m-%d %H:%M:%S"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def get_image_files(folder):
    images = []
    try:
        items = os.listdir(folder)
    except OSError:
        return images
    for item in items:
        if is_image_file(item):
            images.append(item)
    return images


def count_images_recursive(main_folder):
    total = 0
    for _ctime, _folder, images, _rel, _name in list_folders_with_images(main_folder):
        total += len(images)
    return total


def list_folders_with_images(main_folder):
    """列出含图片的文件夹：根目录优先，其余按文件夹创建时间排序。"""
    all_folders = []

    def walk(folder):
        images = get_image_files(folder)
        if images:
            all_folders.append((
                os.path.getctime(folder),
                folder,
                images,
                os.path.relpath(folder, main_folder) if folder != main_folder else "",
                os.path.basename(folder.rstrip("\\/")),
            ))
        try:
            for item in sorted(os.listdir(folder)):
                item_path = os.path.join(folder, item)
                if os.path.isdir(item_path):
                    walk(item_path)
        except OSError:
            return

    walk(main_folder)

    main_entry = None
    rest = []
    for entry in all_folders:
        if os.path.normcase(entry[1]) == os.path.normcase(main_folder):
            main_entry = entry
        else:
            rest.append(entry)
    rest.sort(key=lambda x: x[0])

    result = []
    if main_entry:
        result.append(main_entry)
    result.extend(rest)
    return result

def get_image_timestamp(folder, filename):
    path = os.path.join(folder, filename)
    if platform.system() == "Windows":
        return os.path.getctime(path)
    return os.path.getmtime(path)

def parse_cutoff_datetime(value):
    text = (value or "").strip()
    if not text:
        return None
    for fmt in (DATETIME_FMT, "%Y/%m/%d %H:%M:%S"):
        try:
            return datetime.strptime(text, fmt).timestamp()
        except ValueError:
            continue
    raise ValueError(f"时间格式应为 {DATETIME_FMT}")

def sort_image_files(folder, images, sort_order="asc"):
    reverse = sort_order == "desc"
    return sorted(
        images,
        key=lambda filename: get_image_timestamp(folder, filename),
        reverse=reverse,
    )

def plan_folder_image_pages(folder, images, images_per_page=0, sort_order="asc"):
    if not images:
        return []

    sorted_images = sort_image_files(folder, images, sort_order)
    per_page = len(sorted_images) if images_per_page <= 0 else images_per_page
    per_page = max(1, per_page)
    return [
        sorted_images[index:index + per_page]
        for index in range(0, len(sorted_images), per_page)
    ]

def add_image_pages(
    prs,
    folder,
    images,
    blank_slide_layout,
    slide_w,
    slide_h,
    crop,
    crop_output_dir,
    log_callback,
    images_per_page=0,
    sort_order="asc",
):
    for chunk in plan_folder_image_pages(folder, images, images_per_page, sort_order):
        slide = prs.slides.add_slide(blank_slide_layout)
        image_paths = [os.path.join(folder, filename) for filename in chunk]
        add_images_to_slide(
            slide, image_paths, slide_w, slide_h, crop, crop_output_dir, log_callback
        )

def get_file_prefix(filename, separator="_"):
    name = os.path.splitext(filename)[0]
    if separator and separator in name:
        return name.split(separator, 1)[0]
    return name

def sanitize_filename(name):
    safe_name = re.sub(r'[<>:"/\\|?*]', '_', name).strip()
    return safe_name or "output"

def create_presentation():
    template_candidates = [
        get_resource_path(os.path.join("templates", "default.pptx")),
        get_resource_path("default.pptx"),
    ]
    try:
        import pptx
        template_candidates.append(
            os.path.join(os.path.dirname(pptx.__file__), "templates", "default.pptx")
        )
    except ImportError:
        pass

    for template_path in template_candidates:
        if template_path and os.path.exists(template_path):
            return Presentation(template_path)

    raise FileNotFoundError("未找到 PPT 模板文件 default.pptx")

def group_images_by_prefix(folder, separator="_", cutoff_timestamp=None):
    groups = {}
    for item in os.listdir(folder):
        if is_image_file(item):
            if cutoff_timestamp is not None:
                if get_image_timestamp(folder, item) < cutoff_timestamp:
                    continue
            prefix = get_file_prefix(item, separator)
            if prefix not in groups:
                groups[prefix] = []
            groups[prefix].append(item)
    return groups

def generate_ppt(
    main_folder,
    output_dir,
    mode="folder",
    crop=False,
    crop_output_dir=None,
    log_callback=None,
    prefix_separator="_",
    images_per_page=0,
    sort_order="asc",
    prefix_cutoff_time=None,
):
    prs = create_presentation()
    
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    blank_slide_layout = prs.slide_layouts[6]

    if mode == "prefix":
        cutoff_timestamp = parse_cutoff_datetime(prefix_cutoff_time)
        all_groups = []
        def process_folder(folder):
            rel_path = os.path.relpath(folder, main_folder) if folder != main_folder else ""
            prefix_groups = group_images_by_prefix(
                folder, prefix_separator, cutoff_timestamp
            )
            for prefix, images in prefix_groups.items():
                if not images:
                    continue
                images.sort()
                all_groups.append((rel_path, prefix, folder, images))
            for item in os.listdir(folder):
                item_path = os.path.join(folder, item)
                if os.path.isdir(item_path):
                    process_folder(item_path)
        process_folder(main_folder)
        all_groups.sort(key=lambda x: (x[0], x[1]))
        for rel_path, prefix, folder, images in all_groups:
            slide = prs.slides.add_slide(blank_slide_layout)
            image_paths = [os.path.join(folder, f) for f in images]
            add_images_to_slide(slide, image_paths, slide_w, slide_h, crop, crop_output_dir, log_callback)
    else:
        for _ctime, folder, images, _rel, _name in list_folders_with_images(main_folder):
            add_image_pages(
                prs,
                folder,
                images,
                blank_slide_layout,
                slide_w,
                slide_h,
                crop,
                crop_output_dir,
                log_callback,
                images_per_page,
                sort_order,
            )

    folder_name = sanitize_filename(os.path.basename(main_folder.rstrip("\\/")))
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    os.makedirs(output_dir, exist_ok=True)
    mode_suffix = "按前缀" if mode == "prefix" else "按文件夹"
    output_file = os.path.join(output_dir, f"奢领竞_{folder_name}_{mode_suffix}_{timestamp}.pptx")
    prs.save(output_file)

    return output_file, len(prs.slides)